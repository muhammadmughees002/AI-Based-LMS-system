main.py  —  LMS FastAPI Backend
================================
Endpoints
---------
AUTH
  POST /api/auth/register
  POST /api/auth/login

TEACHER
  POST   /api/teacher/upload-course        (upload PPT/PDF → triggers Ollama generation)
  GET    /api/teacher/courses              (list own courses)
  GET    /api/teacher/course/{id}/students (student submissions for a course)
  POST   /api/teacher/course/{id}/evaluate-all  (trigger Ollama to evaluate all pending)

STUDENT
  GET    /api/student/courses              (list available courses + enrollment status)
  POST   /api/student/enroll/{course_id}
  GET    /api/student/quiz/{course_id}     (get generated quiz for enrolled course)
  POST   /api/student/quiz/{course_id}/submit
  GET    /api/student/assignment/{course_id}
  POST   /api/student/assignment/{course_id}/submit
  GET    /api/student/profile             (all scores + performance history)

Run
---
  pip install fastapi uvicorn python-multipart pypdf python-pptx python-docx requests
  uvicorn main:app --reload --port 8000
"""

import json
import re
import sqlite3
import uuid
from datetime import datetime
from pathlib import Path
from typing import Optional

import requests
from fastapi import (
    BackgroundTasks, Depends, FastAPI, File, Form,
    HTTPException, UploadFile, status,
)
from fastapi.middleware.cors import CORSMiddleware
from fastapi.security import HTTPBasic, HTTPBasicCredentials
from pydantic import BaseModel

# ── optional doc-parsing imports ─────────────────────────────────────────────
try:
    from pypdf import PdfReader
except ImportError:
    PdfReader = None
try:
    from pptx import Presentation
except ImportError:
    Presentation = None
try:
    import docx as _docx
except ImportError:
    _docx = None

# ═══════════════════════════════════════════════════════════════════════════════
# CONFIG
# ═══════════════════════════════════════════════════════════════════════════════

OLLAMA_BASE_URL   = "http://localhost:11434"
OLLAMA_MODEL      = "llama3:8b"
OLLAMA_TIMEOUT    = 600
MAX_CONTENT_CHARS = 6000
DB_PATH           = "lms.db"
UPLOAD_DIR        = Path("uploads")
UPLOAD_DIR.mkdir(exist_ok=True)

app = FastAPI(title="AI-LMS API", version="1.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

security = HTTPBasic()

# ═══════════════════════════════════════════════════════════════════════════════
# DATABASE  (SQLite — swap for PostgreSQL in production via psycopg2)
# ═══════════════════════════════════════════════════════════════════════════════

def get_db():
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    try:
        yield conn
    finally:
        conn.close()


def init_db():
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    c.executescript("""
    CREATE TABLE IF NOT EXISTS users (
        id          TEXT PRIMARY KEY,
        username    TEXT UNIQUE NOT NULL,
        password    TEXT NOT NULL,          -- plain text for demo; hash in prod
        role        TEXT NOT NULL,          -- 'teacher' | 'student'
        created_at  TEXT NOT NULL
    );

    CREATE TABLE IF NOT EXISTS courses (
        id              TEXT PRIMARY KEY,
        teacher_id      TEXT NOT NULL,
        title           TEXT NOT NULL,
        topic           TEXT NOT NULL,
        file_path       TEXT NOT NULL,
        content_preview TEXT,
        quiz_json       TEXT,               -- generated quiz (JSON string)
        assignment_json TEXT,               -- generated assignment (JSON string)
        status          TEXT DEFAULT 'processing',  -- processing | ready | error
        created_at      TEXT NOT NULL,
        FOREIGN KEY (teacher_id) REFERENCES users(id)
    );

    CREATE TABLE IF NOT EXISTS enrollments (
        id          TEXT PRIMARY KEY,
        student_id  TEXT NOT NULL,
        course_id   TEXT NOT NULL,
        enrolled_at TEXT NOT NULL,
        UNIQUE(student_id, course_id),
        FOREIGN KEY (student_id) REFERENCES users(id),
        FOREIGN KEY (course_id)  REFERENCES courses(id)
    );

    CREATE TABLE IF NOT EXISTS quiz_submissions (
        id              TEXT PRIMARY KEY,
        student_id      TEXT NOT NULL,
        course_id       TEXT NOT NULL,
        answers_json    TEXT NOT NULL,      -- {question_id: chosen_option}
        score           REAL,
        total           REAL,
        percentage      REAL,
        feedback_json   TEXT,               -- per-question feedback
        submitted_at    TEXT NOT NULL,
        FOREIGN KEY (student_id) REFERENCES users(id),
        FOREIGN KEY (course_id)  REFERENCES courses(id)
    );

    CREATE TABLE IF NOT EXISTS assignment_submissions (
        id              TEXT PRIMARY KEY,
        student_id      TEXT NOT NULL,
        course_id       TEXT NOT NULL,
        answers_json    TEXT NOT NULL,      -- {task_number: answer_text}
        score           REAL,
        total           REAL,
        percentage      REAL,
        feedback_json   TEXT,
        evaluated       INTEGER DEFAULT 0,
        submitted_at    TEXT NOT NULL,
        FOREIGN KEY (student_id) REFERENCES users(id),
        FOREIGN KEY (course_id)  REFERENCES courses(id)
    );
    """)
    conn.commit()
    conn.close()


init_db()

# ═══════════════════════════════════════════════════════════════════════════════
# HELPERS — AUTH
# ═══════════════════════════════════════════════════════════════════════════════

def get_user(username: str, db: sqlite3.Connection) -> Optional[dict]:
    row = db.execute("SELECT * FROM users WHERE username=?", (username,)).fetchone()
    return dict(row) if row else None


def require_auth(credentials: HTTPBasicCredentials = Depends(security),
                 db: sqlite3.Connection = Depends(get_db)):
    user = get_user(credentials.username, db)
    if not user or user["password"] != credentials.password:
        raise HTTPException(status_code=401, detail="Invalid credentials")
    return user


def require_teacher(user=Depends(require_auth)):
    if user["role"] != "teacher":
        raise HTTPException(status_code=403, detail="Teacher access required")
    return user


def require_student(user=Depends(require_auth)):
    if user["role"] != "student":
        raise HTTPException(status_code=403, detail="Student access required")
    return user

# ═══════════════════════════════════════════════════════════════════════════════
# HELPERS — DOCUMENT PARSING
# ═══════════════════════════════════════════════════════════════════════════════

def extract_text(file_path: str) -> str:
    path = Path(file_path)
    ext  = path.suffix.lower()
    text = ""

    if ext == ".pdf":
        if PdfReader is None:
            raise RuntimeError("pypdf not installed")
        reader = PdfReader(file_path)
        text = "\n\n".join(p.extract_text() or "" for p in reader.pages)

    elif ext == ".pptx":
        if Presentation is None:
            raise RuntimeError("python-pptx not installed")
        prs = Presentation(file_path)
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            parts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = " ".join(r.text for r in para.runs).strip()
                        if line:
                            parts.append(line)
            if parts:
                slides.append(f"[Slide {i}]\n" + "\n".join(parts))
        text = "\n\n".join(slides)

    elif ext in (".docx", ".doc"):
        if _docx is None:
            raise RuntimeError("python-docx not installed")
        document = _docx.Document(file_path)
        text = "\n".join(p.text for p in document.paragraphs if p.text.strip())

    elif ext in (".txt", ".md"):
        text = path.read_text(encoding="utf-8", errors="ignore")

    else:
        raise ValueError(f"Unsupported file type: {ext}")

    if len(text) > MAX_CONTENT_CHARS:
        text = text[:MAX_CONTENT_CHARS] + "\n[...truncated...]"
    return text.strip()

# ═══════════════════════════════════════════════════════════════════════════════
# HELPERS — OLLAMA
# ═══════════════════════════════════════════════════════════════════════════════

def ollama_generate(prompt: str, system: str = "") -> str:
    payload = {
        "model":  OLLAMA_MODEL,
        "prompt": prompt,
        "stream": True,
        "options": {"num_predict": 2048, "temperature": 0.3},
    }
    if system:
        payload["system"] = system

    resp = requests.post(
        f"{OLLAMA_BASE_URL}/api/generate",
        json=payload,
        timeout=OLLAMA_TIMEOUT,
        stream=True,
    )
    resp.raise_for_status()

    tokens = []
    for line in resp.iter_lines():
        if not line:
            continue
        try:
            chunk = json.loads(line)
        except json.JSONDecodeError:
            continue
        tokens.append(chunk.get("response", ""))
        if chunk.get("done"):
            break
    return "".join(tokens).strip()


def parse_json_response(raw: str) -> dict | list:
    """Extract JSON from model response (handles fences, surrounding prose)."""
    try:
        return json.loads(raw)
    except json.JSONDecodeError:
        pass
    fenced = re.search(r"```(?:json)?\s*([\s\S]+?)\s*```", raw)
    if fenced:
        try:
            return json.loads(fenced.group(1))
        except json.JSONDecodeError:
            pass
    for sc, ec in [('{', '}'), ('[', ']')]:
        idx = raw.find(sc)
        if idx != -1:
            depth = 0
            for i, ch in enumerate(raw[idx:], idx):
                if ch == sc:
                    depth += 1
                elif ch == ec:
                    depth -= 1
                    if depth == 0:
                        try:
                            return json.loads(raw[idx:i+1])
                        except json.JSONDecodeError:
                            break
    raise ValueError(f"Cannot parse JSON:\n{raw[:400]}")


def generate_quiz(content: str, topic: str, difficulty: str, num_q: int = 3) -> dict:
    system = ("You are an educational assessment AI. Respond with valid JSON only. "
              "No prose, no markdown fences.")
    prompt = f"""Generate a {difficulty.upper()} difficulty MCQ quiz about "{topic}".

Course content:
---
{content}
---

Rules: exactly {num_q} questions, 4 options (A-D), keep explanations to 1 sentence.

Return ONLY this JSON:
{{
  "topic": "{topic}",
  "difficulty": "{difficulty}",
  "questions": [
    {{
      "id": 1,
      "question": "...",
      "options": {{"A": "...", "B": "...", "C": "...", "D": "..."}},
      "correct_answer": "A",
      "explanation": "..."
    }}
  ]
}}"""
    raw = ollama_generate(prompt, system)
    return parse_json_response(raw)


def generate_assignment(content: str, topic: str, difficulty: str) -> dict:
    system = ("You are a curriculum designer AI. Respond with valid JSON only. "
              "No prose, no markdown fences.")
    prompt = f"""Generate a {difficulty.upper()} difficulty assignment about "{topic}".

Course content:
---
{content}
---

Rules: 2-3 tasks, total_marks=100, keep instructions concise.

Return ONLY this JSON:
{{
  "topic": "{topic}",
  "difficulty": "{difficulty}",
  "title": "...",
  "description": "...",
  "tasks": [
    {{"task_number": 1, "instructions": "...", "marks": 40}},
    {{"task_number": 2, "instructions": "...", "marks": 60}}
  ],
  "total_marks": 100
}}"""
    raw = ollama_generate(prompt, system)
    return parse_json_response(raw)


def evaluate_quiz_submission(quiz: dict, answers: dict) -> dict:
    """Grade MCQ locally (no Ollama needed — answers are in the quiz JSON)."""
    questions  = quiz.get("questions", [])
    score      = 0
    total      = len(questions)
    feedback   = []

    for q in questions:
        qid     = str(q["id"])
        correct = q["correct_answer"]
        chosen  = answers.get(qid, "").upper()
        is_correct = chosen == correct
        if is_correct:
            score += 1
        feedback.append({
            "question_id":  qid,
            "question":     q["question"],
            "your_answer":  chosen,
            "correct":      correct,
            "is_correct":   is_correct,
            "explanation":  q.get("explanation", ""),
        })

    percentage = round((score / total) * 100, 1) if total else 0
    return {
        "score":      score,
        "total":      total,
        "percentage": percentage,
        "feedback":   feedback,
    }


def evaluate_assignment_with_ollama(assignment: dict, answers: dict) -> dict:
    """Use Ollama to evaluate written assignment tasks."""
    system = "You are a fair academic evaluator. Respond with valid JSON only."
    tasks    = assignment.get("tasks", [])
    results  = []
    total_score = 0
    total_marks = 0

    for task in tasks:
        tn      = str(task["task_number"])
        answer  = answers.get(tn, "(no answer provided)")
        marks   = task.get("marks", 10)
        total_marks += marks

        prompt = f"""Evaluate this student answer.

Task: {task['instructions']}
Max marks: {marks}
Student answer: {answer}

Return ONLY this JSON:
{{
  "task_number": {tn},
  "score": <int 0-{marks}>,
  "max_marks": {marks},
  "feedback": "2 sentence feedback",
  "strengths": ["point"],
  "improvements": ["point"]
}}"""
        try:
            raw    = ollama_generate(prompt, system)
            result = parse_json_response(raw)
            total_score += result.get("score", 0)
            results.append(result)
        except Exception as e:
            results.append({
                "task_number":  int(tn),
                "score":        0,
                "max_marks":    marks,
                "feedback":     f"Evaluation error: {str(e)}",
                "strengths":    [],
                "improvements": [],
            })

    percentage = round((total_score / total_marks) * 100, 1) if total_marks else 0
    return {
        "score":      total_score,
        "total":      total_marks,
        "percentage": percentage,
        "tasks":      results,
    }


def compute_difficulty(avg_score: float) -> str:
    if avg_score < 40:
        return "low"
    elif avg_score <= 70:
        return "medium"
    return "hard"

# ═══════════════════════════════════════════════════════════════════════════════
# BACKGROUND TASK — called after file upload
# ═══════════════════════════════════════════════════════════════════════════════

def process_course_background(course_id: str, file_path: str, topic: str):
    """Extract text, generate quiz + assignment, save to DB."""
    conn = sqlite3.connect(DB_PATH)
    try:
        content = extract_text(file_path)
        preview = content[:300]

        quiz       = generate_quiz(content, topic, difficulty="medium")
        assignment = generate_assignment(content, topic, difficulty="medium")

        conn.execute(
            """UPDATE courses SET
               content_preview=?, quiz_json=?, assignment_json=?, status='ready'
               WHERE id=?""",
            (preview, json.dumps(quiz), json.dumps(assignment), course_id)
        )
        conn.commit()
    except Exception as e:
        conn.execute(
            "UPDATE courses SET status=? WHERE id=?",
            (f"error: {str(e)[:200]}", course_id)
        )
        conn.commit()
    finally:
        conn.close()

# ═══════════════════════════════════════════════════════════════════════════════
# PYDANTIC MODELS
# ═══════════════════════════════════════════════════════════════════════════════

class RegisterRequest(BaseModel):
    username: str
    password: str
    role:     str   # 'teacher' | 'student'

class LoginRequest(BaseModel):
    username: str
    password: str

class QuizSubmitRequest(BaseModel):
    answers: dict   # {"1": "A", "2": "C", ...}

class AssignmentSubmitRequest(BaseModel):
    answers: dict   # {"1": "answer text", "2": "answer text", ...}

# ═══════════════════════════════════════════════════════════════════════════════
# AUTH ROUTES
# ═══════════════════════════════════════════════════════════════════════════════

@app.post("/api/auth/register")
def register(req: RegisterRequest, db: sqlite3.Connection = Depends(get_db)):
    if req.role not in ("teacher", "student"):
        raise HTTPException(400, "role must be 'teacher' or 'student'")
    existing = db.execute("SELECT id FROM users WHERE username=?", (req.username,)).fetchone()
    if existing:
        raise HTTPException(400, "Username already taken")
    uid = str(uuid.uuid4())
    db.execute(
        "INSERT INTO users VALUES (?,?,?,?,?)",
        (uid, req.username, req.password, req.role, datetime.utcnow().isoformat())
    )
    db.commit()
    return {"message": "Registered successfully", "user_id": uid, "role": req.role}


@app.post("/api/auth/login")
def login(req: LoginRequest, db: sqlite3.Connection = Depends(get_db)):
    user = get_user(req.username, db)
    if not user or user["password"] != req.password:
        raise HTTPException(401, "Invalid credentials")
    return {
        "user_id":  user["id"],
        "username": user["username"],
        "role":     user["role"],
    }

# ═══════════════════════════════════════════════════════════════════════════════
# TEACHER ROUTES
# ═══════════════════════════════════════════════════════════════════════════════

@app.post("/api/teacher/upload-course")
async def upload_course(
    background_tasks: BackgroundTasks,
    title: str = Form(...),
    topic: str = Form(...),
    file:  UploadFile = File(...),
    user=Depends(require_teacher),
    db: sqlite3.Connection = Depends(get_db),
):
    # Save uploaded file
    ext       = Path(file.filename).suffix
    course_id = str(uuid.uuid4())
    file_path = UPLOAD_DIR / f"{course_id}{ext}"
    content   = await file.read()
    file_path.write_bytes(content)

    # Insert course row (status=processing)
    db.execute(
        "INSERT INTO courses VALUES (?,?,?,?,?,?,?,?,?,?)",
        (course_id, user["id"], title, topic, str(file_path),
         None, None, None, "processing", datetime.utcnow().isoformat())
    )
    db.commit()

    # Kick off background generation
    background_tasks.add_task(
        process_course_background, course_id, str(file_path), topic
    )

    return {
        "message":   "Course uploaded. AI is generating quiz & assignment in the background.",
        "course_id": course_id,
        "status":    "processing",
    }


@app.get("/api/teacher/courses")
def teacher_courses(user=Depends(require_teacher), db: sqlite3.Connection = Depends(get_db)):
    rows = db.execute(
        "SELECT id, title, topic, status, created_at FROM courses WHERE teacher_id=?",
        (user["id"],)
    ).fetchall()
    return [dict(r) for r in rows]


@app.get("/api/teacher/course/{course_id}")
def teacher_course_detail(
    course_id: str,
    user=Depends(require_teacher),
    db: sqlite3.Connection = Depends(get_db),
):
    row = db.execute("SELECT * FROM courses WHERE id=? AND teacher_id=?",
                     (course_id, user["id"])).fetchone()
    if not row:
        raise HTTPException(404, "Course not found")
    c = dict(row)
    c["quiz"]       = json.loads(c["quiz_json"])       if c["quiz_json"]       else None
    c["assignment"] = json.loads(c["assignment_json"]) if c["assignment_json"] else None
    del c["quiz_json"], c["assignment_json"]
    return c


@app.get("/api/teacher/course/{course_id}/students")
def course_students(
    course_id: str,
    user=Depends(require_teacher),
    db: sqlite3.Connection = Depends(get_db),
):
    # Verify ownership
    owned = db.execute("SELECT id FROM courses WHERE id=? AND teacher_id=?",
                       (course_id, user["id"])).fetchone()
    if not owned:
        raise HTTPException(404, "Course not found")

    students = db.execute("""
        SELECT u.username,
               qs.score  AS quiz_score,  qs.total  AS quiz_total,
               qs.percentage AS quiz_pct, qs.submitted_at AS quiz_submitted,
               asub.score AS assign_score, asub.total AS assign_total,
               asub.percentage AS assign_pct, asub.evaluated AS assign_evaluated
        FROM enrollments e
        JOIN users u ON u.id = e.student_id
        LEFT JOIN quiz_submissions qs
               ON qs.student_id = e.student_id AND qs.course_id = e.course_id
        LEFT JOIN assignment_submissions asub
               ON asub.student_id = e.student_id AND asub.course_id = e.course_id
        WHERE e.course_id = ?
    """, (course_id,)).fetchall()

    return [dict(s) for s in students]


@app.post("/api/teacher/course/{course_id}/evaluate-assignments")
def evaluate_all_assignments(
    course_id: str,
    background_tasks: BackgroundTasks,
    user=Depends(require_teacher),
    db: sqlite3.Connection = Depends(get_db),
):
    """Trigger Ollama evaluation for all un-evaluated assignment submissions."""
    owned = db.execute("SELECT assignment_json FROM courses WHERE id=? AND teacher_id=?",
                       (course_id, user["id"])).fetchone()
    if not owned:
        raise HTTPException(404, "Course not found")

    pending = db.execute(
        "SELECT id, answers_json FROM assignment_submissions WHERE course_id=? AND evaluated=0",
        (course_id,)
    ).fetchall()

    assignment = json.loads(owned["assignment_json"]) if owned["assignment_json"] else {}

    def _evaluate_all():
        for sub in pending:
            answers = json.loads(sub["answers_json"])
            result  = evaluate_assignment_with_ollama(assignment, answers)
            conn2   = sqlite3.connect(DB_PATH)
            conn2.execute("""
                UPDATE assignment_submissions
                SET score=?, total=?, percentage=?, feedback_json=?, evaluated=1
                WHERE id=?
            """, (result["score"], result["total"], result["percentage"],
                  json.dumps(result["tasks"]), sub["id"]))
            conn2.commit()
            conn2.close()

    background_tasks.add_task(_evaluate_all)
    return {"message": f"Evaluating {len(pending)} submissions in background."}

# ═══════════════════════════════════════════════════════════════════════════════
# STUDENT ROUTES
# ═══════════════════════════════════════════════════════════════════════════════

@app.get("/api/student/courses")
def list_courses(user=Depends(require_student), db: sqlite3.Connection = Depends(get_db)):
    rows = db.execute("""
        SELECT c.id, c.title, c.topic, c.status, c.created_at,
               CASE WHEN e.id IS NOT NULL THEN 1 ELSE 0 END AS enrolled
        FROM courses c
        LEFT JOIN enrollments e ON e.course_id=c.id AND e.student_id=?
        WHERE c.status='ready'
    """, (user["id"],)).fetchall()
    return [dict(r) for r in rows]


@app.post("/api/student/enroll/{course_id}")
def enroll(course_id: str, user=Depends(require_student),
           db: sqlite3.Connection = Depends(get_db)):
    exists = db.execute("SELECT id FROM courses WHERE id=? AND status='ready'",
                        (course_id,)).fetchone()
    if not exists:
        raise HTTPException(404, "Course not found or not ready")
    try:
        db.execute(
            "INSERT INTO enrollments VALUES (?,?,?,?)",
            (str(uuid.uuid4()), user["id"], course_id, datetime.utcnow().isoformat())
        )
        db.commit()
    except sqlite3.IntegrityError:
        raise HTTPException(400, "Already enrolled")
    return {"message": "Enrolled successfully"}


def _check_enrolled(student_id, course_id, db):
    e = db.execute(
        "SELECT id FROM enrollments WHERE student_id=? AND course_id=?",
        (student_id, course_id)
    ).fetchone()
    if not e:
        raise HTTPException(403, "Not enrolled in this course")


@app.get("/api/student/quiz/{course_id}")
def get_quiz(course_id: str, user=Depends(require_student),
             db: sqlite3.Connection = Depends(get_db)):
    _check_enrolled(user["id"], course_id, db)
    row = db.execute("SELECT quiz_json, title, topic FROM courses WHERE id=?",
                     (course_id,)).fetchone()
    if not row or not row["quiz_json"]:
        raise HTTPException(404, "Quiz not ready yet")
    quiz = json.loads(row["quiz_json"])
    # Strip correct answers before sending to student
    for q in quiz.get("questions", []):
        q.pop("correct_answer", None)
        q.pop("explanation", None)
    return {"course_title": row["title"], "quiz": quiz}


@app.post("/api/student/quiz/{course_id}/submit")
def submit_quiz(
    course_id: str,
    req: QuizSubmitRequest,
    user=Depends(require_student),
    db: sqlite3.Connection = Depends(get_db),
):
    _check_enrolled(user["id"], course_id, db)

    # Check if already submitted
    existing = db.execute(
        "SELECT id FROM quiz_submissions WHERE student_id=? AND course_id=?",
        (user["id"], course_id)
    ).fetchone()
    if existing:
        raise HTTPException(400, "Quiz already submitted")

    row = db.execute("SELECT quiz_json FROM courses WHERE id=?", (course_id,)).fetchone()
    if not row or not row["quiz_json"]:
        raise HTTPException(404, "Quiz not found")

    quiz   = json.loads(row["quiz_json"])
    result = evaluate_quiz_submission(quiz, req.answers)

    db.execute(
        "INSERT INTO quiz_submissions VALUES (?,?,?,?,?,?,?,?,?)",
        (str(uuid.uuid4()), user["id"], course_id,
         json.dumps(req.answers), result["score"], result["total"],
         result["percentage"], json.dumps(result["feedback"]),
         datetime.utcnow().isoformat())
    )
    db.commit()
    return result


@app.get("/api/student/assignment/{course_id}")
def get_assignment(course_id: str, user=Depends(require_student),
                   db: sqlite3.Connection = Depends(get_db)):
    _check_enrolled(user["id"], course_id, db)
    row = db.execute("SELECT assignment_json, title FROM courses WHERE id=?",
                     (course_id,)).fetchone()
    if not row or not row["assignment_json"]:
        raise HTTPException(404, "Assignment not ready yet")
    return {"course_title": row["title"], "assignment": json.loads(row["assignment_json"])}


@app.post("/api/student/assignment/{course_id}/submit")
def submit_assignment(
    course_id: str,
    req: AssignmentSubmitRequest,
    user=Depends(require_student),
    db: sqlite3.Connection = Depends(get_db),
):
    _check_enrolled(user["id"], course_id, db)

    existing = db.execute(
        "SELECT id FROM assignment_submissions WHERE student_id=? AND course_id=?",
        (user["id"], course_id)
    ).fetchone()
    if existing:
        raise HTTPException(400, "Assignment already submitted")

    db.execute(
        "INSERT INTO assignment_submissions VALUES (?,?,?,?,?,?,?,?,?,?)",
        (str(uuid.uuid4()), user["id"], course_id,
         json.dumps(req.answers), None, None, None, None, 0,
         datetime.utcnow().isoformat())
    )
    db.commit()
    return {"message": "Assignment submitted. Teacher will trigger AI evaluation."}


@app.get("/api/student/profile")
def student_profile(user=Depends(require_student), db: sqlite3.Connection = Depends(get_db)):
    quiz_rows = db.execute("""
        SELECT c.title, c.topic, qs.score, qs.total, qs.percentage,
               qs.feedback_json, qs.submitted_at
        FROM quiz_submissions qs
        JOIN courses c ON c.id = qs.course_id
        WHERE qs.student_id=?
        ORDER BY qs.submitted_at DESC
    """, (user["id"],)).fetchall()

    assign_rows = db.execute("""
        SELECT c.title, c.topic, asub.score, asub.total, asub.percentage,
               asub.feedback_json, asub.evaluated, asub.submitted_at
        FROM assignment_submissions asub
        JOIN courses c ON c.id = asub.course_id
        WHERE asub.student_id=?
        ORDER BY asub.submitted_at DESC
    """, (user["id"],)).fetchall()

    quiz_scores = [r["percentage"] for r in quiz_rows if r["percentage"] is not None]
    avg_score   = round(sum(quiz_scores) / len(quiz_scores), 1) if quiz_scores else None
    level       = compute_difficulty(avg_score) if avg_score is not None else "not assessed"

    def parse_feedback(row):
        d = dict(row)
        if d.get("feedback_json"):
            d["feedback"] = json.loads(d["feedback_json"])
        del d["feedback_json"]
        return d

    return {
        "username":    user["username"],
        "level":       level,
        "avg_score":   avg_score,
        "quizzes":     [parse_feedback(r) for r in quiz_rows],
        "assignments": [parse_feedback(r) for r in assign_rows],
    }


@app.get("/api/student/quiz/{course_id}/result")
def quiz_result(course_id: str, user=Depends(require_student),
                db: sqlite3.Connection = Depends(get_db)):
    row = db.execute(
        "SELECT * FROM quiz_submissions WHERE student_id=? AND course_id=?",
        (user["id"], course_id)
    ).fetchone()
    if not row:
        raise HTTPException(404, "No submission found")
    d = dict(row)
    d["feedback"] = json.loads(d["feedback_json"])
    del d["feedback_json"], d["answers_json"]
    return d


@app.get("/api/student/assignment/{course_id}/result")
def assignment_result(course_id: str, user=Depends(require_student),
                      db: sqlite3.Connection = Depends(get_db)):
    row = db.execute(
        "SELECT * FROM assignment_submissions WHERE student_id=? AND course_id=?",
        (user["id"], course_id)
    ).fetchone()
    if not row:
        raise HTTPException(404, "No submission found")
    d = dict(row)
    if d["feedback_json"]:
        d["feedback"] = json.loads(d["feedback_json"])
    del d["feedback_json"], d["answers_json"]
    return d


# ── health check ─────────────────────────────────────────────────────────────
@app.get("/api/health")
def health():
    return {"status": "ok", "ollama": OLLAMA_BASE_URL, "model": OLLAMA_MODEL}